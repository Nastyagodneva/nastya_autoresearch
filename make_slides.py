"""
make_slides.py  —  AutoResearch progress slides
Run:  uv run make_slides.py
Requires: python-pptx   (uv add python-pptx)
Output:   AutoResearch_Progress.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── helpers ───────────────────────────────────────────────────────────────────

W, H = Inches(10), Inches(7.5)   # 4:3 slide

def rgb(hex6):
    h = hex6.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def add_box(slide, left, top, width, height,
            fill_hex=None, border_hex=None):
    """Add a plain rectangle; return the shape."""
    box = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    fill = box.fill
    if fill_hex:
        fill.solid()
        fill.fore_color.rgb = rgb(fill_hex)
    else:
        fill.background()

    line = box.line
    if border_hex:
        line.color.rgb = rgb(border_hex)
        line.width = Pt(0.75)
    else:
        line.fill.background()
    return box

def add_textbox(slide, left, top, width, height):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = True
    return txb

def para(tf, text, size=11, bold=False, color="000000", space_before=0):
    p = tf.add_paragraph()
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.name = "Georgia"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return p

def colored_box(slide, left, top, width, height,
                heading, lines,
                fill_hex, border_hex,
                head_size=11, body_size=10.5):
    """Filled box with optional bold heading + body lines."""
    box = add_box(slide, left, top, width, height, fill_hex, border_hex)
    tf = box.text_frame
    tf.word_wrap = True

    first = True
    if heading:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = heading
        run.font.name = "Georgia"
        run.font.size = Pt(head_size)
        run.font.bold = True
        run.font.color.rgb = rgb("111111")

    for line in lines:
        p = tf.paragraphs[0] if (first and not heading) else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = line
        run.font.name = "Georgia"
        run.font.size = Pt(body_size)
        run.font.bold = False
        run.font.color.rgb = rgb("222222")

def title_box(slide, text, subtitle=""):
    txb = add_textbox(slide, Inches(0.35), Inches(0.2), Inches(9.4), Inches(0.85))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Georgia"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = rgb("000000")
    if subtitle:
        para(tf, subtitle, size=11, color="555555")

def logo(slide):
    txb = add_textbox(slide, Inches(0.3), Inches(7.05), Inches(3), Inches(0.4))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "THE HUMAN  ·  PHENOTYPE  ·  PROJECT"
    run.font.name = "Georgia"
    run.font.size = Pt(6.5)
    run.font.color.rgb = rgb("666666")
    run.font.bold = False

def barcode(slide):
    txb = add_textbox(slide, Inches(9.3), Inches(0.08), Inches(0.65), Inches(0.5))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "| || ||| || | |||"
    run.font.name = "Georgia"
    run.font.size = Pt(7)
    run.font.color.rgb = rgb("000000")

def footer(slide):
    logo(slide)
    barcode(slide)

def new_slide(prs):
    blank_layout = prs.slide_layouts[6]  # completely blank
    slide = prs.slides.add_slide(blank_layout)
    # white background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb("FFFFFF")
    return slide

# ── Slide builders ────────────────────────────────────────────────────────────

def slide1(prs):
    s = new_slide(prs)
    title_box(s, "Agentic ML workflows for multi-omics phenotype prediction")
    colored_box(s, Inches(0.35), Inches(1.1), Inches(4.7), Inches(2.2),
                "Primary goal",
                ["Develop reusable agentic workflows where an agent",
                 "iteratively improves ML pipelines.",
                 "Phenotype prediction is the testbed, not the endpoint."],
                "E8F5E9", "4CAF50")
    colored_box(s, Inches(5.2), Inches(1.1), Inches(4.5), Inches(2.2),
                "Dataset & setup",
                ["Samples: 5,053",
                 "Raw features: 1,028 Olink proteins",
                 "Embeddings: 128 UKBB-trained",
                 "Targets: 52 continuous phenotypes",
                 "Agent: Cursor (-> Claude Code)",
                 "Loop: edit -> eval -> keep/revert"],
                "F5F5F5", "BBBBBB")
    colored_box(s, Inches(0.35), Inches(3.45), Inches(9.35), Inches(0.9),
                None,
                ["The main purpose is to improve agentic ways of working on ML codebases.",
                 "The agent reads local context, makes small testable edits, runs a fixed evaluation harness, and documents keep/revert decisions."],
                "FFF9C4", "F9A825")
    footer(s)

def slide2(prs):
    s = new_slide(prs)
    title_box(s, "AutoResearch: a workflow pattern, not a library")
    colored_box(s, Inches(0.35), Inches(1.1), Inches(4.7), Inches(2.1),
                "What it is NOT",
                ["A hyperparameter search engine",
                 "AutoML over a fixed model space",
                 "Grid search with extra steps",
                 "A library you import"],
                "FFEBEE", "E53935")
    colored_box(s, Inches(5.2), Inches(1.1), Inches(4.5), Inches(2.1),
                "What it IS",
                ["Agent edits pipeline code itself",
                 "Fixed evaluation harness (agent_score)",
                 "Keep / revert discipline via git",
                 "Accumulated knowledge (skills)"],
                "E8F5E9", "4CAF50")
    colored_box(s, Inches(0.35), Inches(3.35), Inches(9.35), Inches(0.85),
                None,
                ["Key novelty: the agent modifies analysis pipeline logic — preprocessing, feature selection,",
                 "model composition — not just parameter values. This is what separates it from grid search."],
                "FFF9C4", "F9A825")
    colored_box(s, Inches(0.35), Inches(4.3), Inches(9.35), Inches(0.6),
                None,
                ["Loop:  train.py (editable)  ->  agent_loop.py (fixed harness)  ->  agent_score  ->  keep / revert  ->  repeat"],
                "F5F5F5", "BBBBBB")
    footer(s)

def slide3(prs):
    s = new_slide(prs)
    title_box(s, "Benchmark results (pre-agentic baseline)",
              subtitle="468 runs · 52 targets x 3 models x 3 feature sets")
    rows = [
        ("raw+embed  lgbm",       0.451),
        ("raw  lgbm",             0.433),
        ("raw+embed  stacking",   0.429),
        ("raw  stacking",         0.416),
        ("raw+embed  elasticnet", 0.412),
        ("embed  lgbm",           0.376),
    ]
    bar_left  = Inches(3.7)
    bar_max_w = Inches(5.0)
    for i, (label, val) in enumerate(rows):
        y = Inches(1.15) + i * Inches(0.72)
        # label
        txb = add_textbox(s, Inches(0.35), y, Inches(3.3), Inches(0.5))
        para(txb.text_frame, label, size=10.5)
        # bar
        bw = int(bar_max_w * val)
        add_box(s, bar_left, y + Inches(0.07), bw, Inches(0.38), "1565C0", None)
        # value
        txb2 = add_textbox(s, bar_left + bw + Inches(0.1), y, Inches(0.7), Inches(0.5))
        para(txb2.text_frame, f"{val:.3f}", size=10.5, bold=True, color="1565C0")

    colored_box(s, Inches(0.35), Inches(5.6), Inches(9.35), Inches(0.75),
                None,
                ["Role in project: benchmark = static exploration of the model space.",
                 "AutoResearch = dynamic optimization on top."],
                "E8F5E9", "4CAF50")
    footer(s)

def slide4(prs):
    s = new_slide(prs)
    title_box(s, "Agentic optimization results",
              subtitle="Per-phenotype sessions with keep / revert discipline")

    col_lefts  = [Inches(0.35), Inches(1.9),  Inches(2.75), Inches(3.55), Inches(4.45), Inches(8.1)]
    col_widths = [Inches(1.5),  Inches(0.82), Inches(0.76), Inches(0.87), Inches(3.6),  Inches(1.55)]
    headers = ["Phenotype", "Baseline", "Best", "Delta", "Champion", "Status"]

    # header row
    for j, h in enumerate(headers):
        box = add_box(s, col_lefts[j], Inches(1.1), col_widths[j], Inches(0.42), "1565C0", None)
        tf = box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = h
        run.font.name = "Georgia"
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = rgb("FFFFFF")

    table_rows = [
        ["HbA1C (BT) — after hyperparams", "0.287", "0.315", "+0.028", "stacking_linear + raw+embed", "improved"],
        ["HbA1C (BT) — after structural",  "0.315", "0.355", "+0.040", "stacking_linear + raw+embed", "improved"],
        ["Liver sound speed",               "0.180", "0.180", "~0",     "stacking_linear + raw",       "plateau"],
        ["BMI",                             "0.653", "0.653", "0",      "elasticnet + raw+embed",       "saturated"],
    ]
    fills = ["E8F5E9", "C8E6C9", "F5F5F5", "FFFFFF"]
    for i, row in enumerate(table_rows):
        y = Inches(1.52) + i * Inches(0.62)
        for j, cell in enumerate(row):
            box = add_box(s, col_lefts[j], y, col_widths[j], Inches(0.58), fills[i], "CCCCCC")
            tf = box.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = cell
            run.font.name = "Georgia"
            run.font.size = Pt(10)
            run.font.bold = (j == 3)
            run.font.color.rgb = rgb("2E7D32") if (j == 3 and cell.startswith("+")) else rgb("111111")

    colored_box(s, Inches(0.35), Inches(3.9), Inches(9.35), Inches(1.2),
                None,
                ["HbA1C structural session (+0.040): ColumnTransformer separating raw/embed preprocessing",
                 "+ split-view stacking + TransformedTargetRegressor.",
                 "These are pipeline-level changes a grid search would not find."],
                "E8F5E9", "4CAF50")
    footer(s)

def slide5(prs):
    s = new_slide(prs)
    title_box(s, "Discovered modeling skills",
              subtitle="Validated knowledge that accumulates across phenotype sessions")
    colored_box(s, Inches(0.35), Inches(1.1), Inches(4.7), Inches(1.9),
                "Prompt (one-time, forgotten)",
                ['"Try SelectKBest on raw features"',
                 "No transfer.",
                 "Each new phenotype starts from zero."],
                "FFEBEE", "E53935")
    colored_box(s, Inches(5.2), Inches(1.1), Inches(4.5), Inches(1.9),
                "Skill (persists as prior)",
                ['"SelectKBest(k=300) on raw helped HbA1C +0.04.',
                 'PCA hurt. k=500 worse."',
                 "Used as prior for every future phenotype."],
                "E8F5E9", "4CAF50")
    colored_box(s, Inches(0.35), Inches(3.15), Inches(9.35), Inches(3.1),
                "Skills extracted so far",
                ["1.  Separate raw / embed preprocessing (ColumnTransformer)",
                 "2.  Split-view stacking — different learners see different feature views",
                 "3.  TransformedTargetRegressor for skewed / weak targets",
                 "4.  Stacking beats standalone for weak phenotypes (R2 < 0.35)",
                 "5.  Embeddings sometimes hurt — always test raw-only baseline",
                 "6.  Strong phenotypes (R2 > 0.6) saturate fast — stop early",
                 "7.  SelectKBest > PCA for Olink features (signal sparse, not low-rank)"],
                "F5F5F5", "BBBBBB")
    footer(s)

def slide6(prs):
    s = new_slide(prs)
    title_box(s, "Roadmap")
    items = [
        (True,  "1. Infrastructure",               "prepare.py, train.py, agent_loop.py, phenotype workspaces"),
        (True,  "2. Benchmark (468 runs)",          "Exhaustive baseline across all targets x models x features"),
        (True,  "3. Hyperparameter sessions",       "Per-phenotype tuning via agent loop"),
        (True,  "4. Structural exploration (HbA1C)","ColumnTransformer, split-view stacking, +0.040 gain"),
        (True,  "5. skills.md",                    "Formalise validated findings as reusable priors"),
        (False, "6. New phenotypes (2-3)",          "Triglycerides, CRP, eGFR — span weak / medium / strong"),
        (False, "7. Skills-guided sessions",        "Agent reads skills.md before starting — test knowledge reuse"),
        (False, "8. Add LGBM back",                 "As expensive expert model where linear stacking fails"),
        (False, "9. Cross-phenotype meta-learning", "Champion patterns as priors for phenotype similarity"),
        (False, "10. Generalise to new X/Y",        "Minimal-input system for new omics datasets"),
    ]
    for i, (done, step, desc) in enumerate(items):
        y = Inches(0.95) + i * Inches(0.615)
        fill   = "E8F5E9" if done else "F5F5F5"
        border = "4CAF50" if done else "CCCCCC"
        tick   = "[x]" if done else "[ ]"
        colored_box(s, Inches(0.35), y, Inches(9.35), Inches(0.55),
                    None, [f"{tick}  {step}  —  {desc}"],
                    fill, border, body_size=10)
    footer(s)

# ── main ──────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide1(prs)
    slide2(prs)
    slide3(prs)
    slide4(prs)
    slide5(prs)
    slide6(prs)

    out = "AutoResearch_Progress.pptx"
    prs.save(out)
    print(f"Saved: {out}")

if __name__ == "__main__":
    main()